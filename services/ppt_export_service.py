# -*- coding: utf-8 -*-
"""
PPT 导出服务
将统计图表截图 + 关键人员能力画像生成为 PowerPoint 文件

依赖：python-pptx
"""
import io
import base64
import hashlib
import json
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


# ─── PPT 尺寸常量（16:9 宽屏）───────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

# ─── 颜色常量 ────────────────────────────────────────────────────────────────
COLOR_PRIMARY    = RGBColor(0x1A, 0x56, 0xDB)   # 蓝色
COLOR_DANGER     = RGBColor(0xDC, 0x35, 0x45)   # 红色（关键人员）
COLOR_SUCCESS    = RGBColor(0x28, 0xA7, 0x45)   # 绿色
COLOR_BG_DARK    = RGBColor(0x1E, 0x29, 0x3B)   # 深色背景（封面）
COLOR_BG_LIGHT   = RGBColor(0xF8, 0xF9, 0xFA)   # 浅灰背景
COLOR_TEXT_MAIN  = RGBColor(0x21, 0x25, 0x29)   # 主文字色
COLOR_TEXT_MUTED = RGBColor(0x6C, 0x75, 0x7D)   # 次要文字


class PPTExportService:
    THEMES = {
        'blue': {
            'bg': RGBColor(0xF4, 0xF6, 0xF8),
            'header_bg': RGBColor(0x1A, 0x56, 0xDB),
            'header_text': RGBColor(0xFF, 0xFF, 0xFF),
            'card_bg': RGBColor(0xFF, 0xFF, 0xFF),
            'accent': RGBColor(0xDC, 0x35, 0x45),
            'text_main': RGBColor(0x33, 0x33, 0x33)
        },
        'dark': {
            'bg': RGBColor(0x18, 0x18, 0x1B),
            'header_bg': RGBColor(0x27, 0x27, 0x2A),
            'header_text': RGBColor(0xF4, 0xF4, 0xF5),
            'card_bg': RGBColor(0x27, 0x27, 0x2A),
            'accent': RGBColor(0x10, 0xB9, 0x81),
            'text_main': RGBColor(0xF9, 0xFA, 0xFB)
        },
        'simple': {
            'bg': RGBColor(0xFF, 0xFF, 0xFF),
            'header_bg': RGBColor(0xF8, 0xF9, 0xFA),
            'header_text': RGBColor(0x21, 0x25, 0x29),
            'card_bg': RGBColor(0xF8, 0xF9, 0xFA),
            'accent': RGBColor(0x1A, 0x56, 0xDB),
            'text_main': RGBColor(0x11, 0x18, 0x27)
        }
    }

    def __init__(self, theme_name='blue'):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self._blank_layout = self.prs.slide_layouts[6]
        self.theme = self.THEMES.get(theme_name, self.THEMES['blue'])

    def generate(self, start_date: str, end_date: str, module_slides: list,
                 nine_grid_image: str, key_persons: list,
                 summary_data: dict = None) -> bytes:
        # 封面
        self._add_cover_slide(start_date, end_date)

        # 执行摘要页
        if summary_data:
            self._add_summary_slide(start_date, end_date, summary_data)

        # 目录页
        self._add_toc_slide(module_slides, nine_grid_image, key_persons)

        # 各统计模块页
        for slide_def in module_slides:
            images = slide_def.get("images", [])
            title = slide_def.get("title", "")
            note = slide_def.get("note", "")

            if not images:
                self._add_text_only_slide(title, note or f"{title}统计数据（本次未包含图表）")
            else:
                total_slides = (len(images) + 1) // 2
                for i in range(0, len(images), 2):
                    page_num = i // 2 + 1
                    page_title = title if total_slides == 1 else f"{title} ({page_num}/{total_slides})"

                    if i + 1 < len(images):
                        self._add_double_image_slide(page_title, images[i], images[i + 1], note)
                    else:
                        self._add_single_image_slide(page_title, images[i], note)

        # 九宫格页
        if nine_grid_image:
            self._add_single_image_slide("人才九宫格", nine_grid_image, "基于三维分的人才视图")

        # 关键人员分页
        for person in key_persons:
            self._add_person_slide(person)

        # 统一添加页脚和页码（跳过封面）
        self._add_footers(start_date, end_date)

        buf = io.BytesIO()
        self.prs.save(buf)
        buf.seek(0)
        return buf.read()

    def _apply_background(self, slide):
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.theme['bg']
        
    def _draw_card(self, slide, left, top, width, height, title=None):
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.theme['card_bg']
        # Remove border or make it very subtle
        shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        
        if title:
            # 内部标题
            self._add_textbox(slide, text=title,
                              left=left + Inches(0.2), top=top + Inches(0.1),
                              width=width - Inches(0.4), height=Inches(0.4),
                              font_size=Pt(14), bold=True,
                              color=self.theme['text_main'],
                              align=PP_ALIGN.CENTER)

    def _add_cover_slide(self, start_date: str, end_date: str):
        slide = self.prs.slides.add_slide(self._blank_layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = COLOR_BG_DARK

        accent = slide.shapes.add_shape(1, 0, 0, Inches(0.35), SLIDE_HEIGHT)
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.theme['accent']
        accent.line.fill.background()

        from pptx.util import Emu
        line_bar = slide.shapes.add_shape(1, Inches(0.7), Inches(3.15), Inches(8.0), Pt(2))
        line_bar.fill.solid()
        line_bar.fill.fore_color.rgb = RGBColor(0x44, 0x55, 0x66)
        line_bar.line.fill.background()

        self._add_textbox(slide, text="人员综合能力报告", left=Inches(0.8), top=Inches(1.8),
                          width=Inches(10.0), height=Inches(1.4), font_size=Pt(44), bold=True,
                          color=RGBColor(0xFF, 0xFF, 0xFF))

        self._add_textbox(slide, text="Personnel Comprehensive Capability Report",
                          left=Inches(0.8), top=Inches(3.0), width=Inches(10.0), height=Inches(0.5),
                          font_size=Pt(14), color=RGBColor(0x88, 0x99, 0xAA))

        self._add_textbox(slide, text=f"统计区间：{start_date}  至  {end_date}",
                          left=Inches(0.8), top=Inches(3.7), width=Inches(10.0), height=Inches(0.55),
                          font_size=Pt(18), color=RGBColor(0x60, 0xC6, 0xFF))
                          
        self._add_textbox(slide, text=datetime.now().strftime("生成时间：%Y-%m-%d %H:%M"),
                          left=Inches(0.8), top=Inches(6.6), width=Inches(10.0), height=Inches(0.4),
                          font_size=Pt(11), color=RGBColor(0x55, 0x66, 0x77))

    # ── 执行摘要页 ────────────────────────────────────────────────────
    def _add_summary_slide(self, start_date: str, end_date: str, summary_data: dict):
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._apply_background(slide)
        self._draw_slide_header(slide, "执行摘要")

        body_top = Inches(1.05)
        total_persons = summary_data.get('total_persons', 0)
        key_count = summary_data.get('key_persons_count', 0)
        chart_counts = summary_data.get('module_chart_counts', {})
        total_charts = sum(chart_counts.values()) if chart_counts else 0

        # KPI 卡片行
        kpi_data = [
            ('📊', '统计图表', f'{total_charts} 张', self.theme['header_bg']),
            ('👥', '纳入分析人数', f'{total_persons} 人', RGBColor(0x00, 0xB4, 0x2A)),
            ('⚠️', '关键人员', f'{key_count} 人', RGBColor(0xF5, 0x3F, 0x3F)),
            ('📅', '统计区间', f'{start_date} ~ {end_date}', RGBColor(0xFF, 0x7D, 0x00)),
        ]

        card_w = Inches(2.95)
        card_h = Inches(1.3)
        gap = Inches(0.15)
        x_start = Inches(0.2)

        for idx, (icon, label, value, accent_color) in enumerate(kpi_data):
            x = x_start + idx * (card_w + gap)
            # 卡片背景
            card = slide.shapes.add_shape(1, x, body_top, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = self.theme['card_bg']
            card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            # 顶部色条
            accent_bar = slide.shapes.add_shape(1, x, body_top, card_w, Pt(4))
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = accent_color
            accent_bar.line.fill.background()
            # 图标+标签
            self._add_textbox(slide, text=f'{icon} {label}',
                              left=x + Inches(0.15), top=body_top + Inches(0.2),
                              width=card_w - Inches(0.3), height=Inches(0.35),
                              font_size=Pt(11), color=COLOR_TEXT_MUTED)
            # 数值
            self._add_textbox(slide, text=value,
                              left=x + Inches(0.15), top=body_top + Inches(0.6),
                              width=card_w - Inches(0.3), height=Inches(0.5),
                              font_size=Pt(22), bold=True, color=self.theme['text_main'])

        # 模块图表明细
        detail_top = body_top + card_h + Inches(0.35)
        self._draw_card(slide, Inches(0.2), detail_top,
                       SLIDE_WIDTH - Inches(0.4), Inches(3.8), title='各模块图表统计')

        y_detail = detail_top + Inches(0.55)
        for mod_name, count in chart_counts.items():
            self._add_textbox(slide, text=f'▸ {mod_name}',
                              left=Inches(0.5), top=y_detail,
                              width=Inches(6), height=Inches(0.35),
                              font_size=Pt(13), color=self.theme['text_main'])
            self._add_textbox(slide, text=f'{count} 张图表',
                              left=Inches(7), top=y_detail,
                              width=Inches(3), height=Inches(0.35),
                              font_size=Pt(13), bold=True, color=self.theme['header_bg'],
                              align=PP_ALIGN.LEFT)
            y_detail += Inches(0.45)

    # ── 目录页 ────────────────────────────────────────────────────
    def _add_toc_slide(self, module_slides: list, nine_grid_image: str, key_persons: list):
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._apply_background(slide)
        self._draw_slide_header(slide, "目录")

        toc_items = []
        page_offset = 4  # 封面(1) + 摘要(2) + 目录(3)，内容从第4页开始

        for slide_def in module_slides:
            images = slide_def.get('images', [])
            title = slide_def.get('title', '')
            toc_items.append((title, f'{len(images)} 张图表'))
            page_offset += max(1, (len(images) + 1) // 2)

        if nine_grid_image:
            toc_items.append(('人才九宫格', '1 张图表'))
            page_offset += 1

        if key_persons:
            toc_items.append(('关键人员详情', f'{len(key_persons)} 人'))

        y = Inches(1.2)
        for idx, (title, detail) in enumerate(toc_items):
            num = f'{idx + 1:02d}'
            # 序号圆点
            self._add_textbox(slide, text=num,
                              left=Inches(0.6), top=y,
                              width=Inches(0.6), height=Inches(0.5),
                              font_size=Pt(16), bold=True,
                              color=self.theme['header_bg'],
                              align=PP_ALIGN.CENTER)
            # 标题
            self._add_textbox(slide, text=title,
                              left=Inches(1.4), top=y,
                              width=Inches(6), height=Inches(0.5),
                              font_size=Pt(16), bold=False,
                              color=self.theme['text_main'])
            # 详情
            self._add_textbox(slide, text=detail,
                              left=Inches(8), top=y,
                              width=Inches(3.5), height=Inches(0.5),
                              font_size=Pt(12),
                              color=COLOR_TEXT_MUTED,
                              align=PP_ALIGN.LEFT)
            # 分隔线
            sep = slide.shapes.add_shape(1, Inches(1.4), y + Inches(0.5),
                                         Inches(10), Pt(1))
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
            sep.line.fill.background()
            y += Inches(0.7)

    # ── 页脚和页码 ────────────────────────────────────────────────────
    def _add_footers(self, start_date: str, end_date: str):
        total = len(self.prs.slides)
        for idx, slide in enumerate(self.prs.slides):
            if idx == 0:  # 跳过封面
                continue
            # 页码（右下角）
            self._add_textbox(slide, text=f'{idx + 1} / {total}',
                              left=SLIDE_WIDTH - Inches(1.2), top=SLIDE_HEIGHT - Inches(0.4),
                              width=Inches(1.0), height=Inches(0.3),
                              font_size=Pt(9), color=COLOR_TEXT_MUTED,
                              align=PP_ALIGN.RIGHT)
            # 左下角标记
            self._add_textbox(slide, text=f'人员综合能力报告 | {start_date} ~ {end_date}',
                              left=Inches(0.3), top=SLIDE_HEIGHT - Inches(0.4),
                              width=Inches(6.0), height=Inches(0.3),
                              font_size=Pt(8), color=COLOR_TEXT_MUTED)

    def _add_single_image_slide(self, title: str, chart_obj, note: str = ""):
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._apply_background(slide)
        self._draw_slide_header(slide, title)
        
        img_b64 = chart_obj.get("image", "") if isinstance(chart_obj, dict) else chart_obj
        chart_title = chart_obj.get("title", "") if isinstance(chart_obj, dict) else ""

        top_offset = Inches(1.1)
        bottom_margin = Inches(0.8) if note else Inches(0.3)
        card_h = SLIDE_HEIGHT - top_offset - bottom_margin
        card_w = SLIDE_WIDTH - Inches(0.8)
        card_left = Inches(0.4)
        
        self._draw_card(slide, card_left, top_offset, card_w, card_h, chart_title)
        
        img_top = top_offset + (Inches(0.5) if chart_title else Inches(0.2))
        img_h = card_h - (Inches(0.7) if chart_title else Inches(0.4))
        
        self._insert_image(slide, img_b64, card_left, img_top, card_w, img_h)

        if note:
            self._add_textbox(slide, text=note, left=Inches(0.4), top=SLIDE_HEIGHT - Inches(0.6),
                              width=card_w, height=Inches(0.4), font_size=Pt(10), color=self.theme['text_main'])

    def _add_double_image_slide(self, title: str, chart_obj_left, chart_obj_right, note: str = ""):
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._apply_background(slide)
        self._draw_slide_header(slide, title)

        top_offset = Inches(1.1)
        bottom_margin = Inches(0.8) if note else Inches(0.3)
        card_h = SLIDE_HEIGHT - top_offset - bottom_margin
        card_w = (SLIDE_WIDTH - Inches(1.0)) / 2
        gap = Inches(0.2)
        
        # Left
        img_b64_L = chart_obj_left.get("image", "") if isinstance(chart_obj_left, dict) else chart_obj_left
        title_L = chart_obj_left.get("title", "") if isinstance(chart_obj_left, dict) else ""
        self._draw_card(slide, Inches(0.4), top_offset, card_w, card_h, title_L)
        self._insert_image(slide, img_b64_L, Inches(0.4), top_offset + (Inches(0.5) if title_L else Inches(0.2)), 
                           card_w, card_h - (Inches(0.7) if title_L else Inches(0.4)))
                           
        # Right
        img_b64_R = chart_obj_right.get("image", "") if isinstance(chart_obj_right, dict) else chart_obj_right
        title_R = chart_obj_right.get("title", "") if isinstance(chart_obj_right, dict) else ""
        left_R = Inches(0.4) + card_w + gap
        self._draw_card(slide, left_R, top_offset, card_w, card_h, title_R)
        self._insert_image(slide, img_b64_R, left_R, top_offset + (Inches(0.5) if title_R else Inches(0.2)), 
                           card_w, card_h - (Inches(0.7) if title_R else Inches(0.4)))

        if note:
            self._add_textbox(slide, text=note, left=Inches(0.4), top=SLIDE_HEIGHT - Inches(0.6),
                              width=SLIDE_WIDTH - Inches(0.8), height=Inches(0.4), 
                              font_size=Pt(10), color=self.theme['text_main'])

    def _add_person_slide(self, person: dict):
        """
        关键人员详情页（美化三栏布局）

        ┌──────────────── 顶部横幅（姓名 / 工号 / 部门）────────────────┐
        │ 左栏（深色信息卡）│  中栏（雷达图）  │  右栏（评分+诊断）     │
        └───────────────────────────────────────────────────────────────┘
        """
        slide = self.prs.slides.add_slide(self._blank_layout)

        # ── 浅灰色全页背景 ──────────────────────────────────────────────
        bg_fill = slide.background.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF8)

        name       = person.get("name", "")
        emp_no     = person.get("emp_no", "")
        department = person.get("department", "")
        scores     = person.get("scores", {})

        # ── 顶部横幅（深红渐变感）──────────────────────────────────────
        banner_h = Inches(0.9)
        banner = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, banner_h)
        banner.fill.solid()
        banner.fill.fore_color.rgb = RGBColor(0xC0, 0x20, 0x2F)
        banner.line.fill.background()

        # 左侧深色边条（装饰）
        edge = slide.shapes.add_shape(1, 0, 0, Inches(0.28), banner_h)
        edge.fill.solid()
        edge.fill.fore_color.rgb = RGBColor(0x8B, 0x00, 0x0D)
        edge.line.fill.background()

        # 警示图标 + 姓名 + 工号
        self._add_textbox(slide,
                          text=f"⚠  关键人员",
                          left=Inches(0.45), top=Inches(0.05),
                          width=Inches(2.0), height=Inches(0.42),
                          font_size=Pt(11), bold=False,
                          color=RGBColor(0xFF, 0xCC, 0xCC),
                          align=PP_ALIGN.LEFT)

        self._add_textbox(slide,
                          text=f"{name}  （{emp_no}）",
                          left=Inches(0.45), top=Inches(0.42),
                          width=Inches(6.0), height=Inches(0.45),
                          font_size=Pt(20), bold=True,
                          color=RGBColor(0xFF, 0xFF, 0xFF),
                          align=PP_ALIGN.LEFT)

        self._add_textbox(slide,
                          text=department,
                          left=Inches(6.6), top=Inches(0.28),
                          width=Inches(6.4), height=Inches(0.45),
                          font_size=Pt(14),
                          color=RGBColor(0xFF, 0xDD, 0xDD),
                          align=PP_ALIGN.RIGHT)

        body_top = banner_h + Inches(0.15)
        body_h   = SLIDE_HEIGHT - body_top - Inches(0.15)

        # ═══════════════════════════════════════════════════════════════
        # 左栏：深色信息卡（综合评分 + 五维数字）
        # ═══════════════════════════════════════════════════════════════
        left_w  = Inches(2.4)
        left_bg = slide.shapes.add_shape(1, Inches(0.12), body_top, left_w, body_h)
        left_bg.fill.solid()
        left_bg.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        left_bg.line.fill.background()

        # 综合评分大字
        comp_val = scores.get("comprehensive")
        comp_str = f"{comp_val:.0f}" if comp_val is not None else "--"
        self._add_textbox(slide,
                          text="综合评分",
                          left=Inches(0.18), top=body_top + Inches(0.18),
                          width=left_w - Inches(0.12), height=Inches(0.4),
                          font_size=Pt(11),
                          color=RGBColor(0x88, 0xAA, 0xCC),
                          align=PP_ALIGN.CENTER)
        self._add_textbox(slide,
                          text=comp_str,
                          left=Inches(0.18), top=body_top + Inches(0.5),
                          width=left_w - Inches(0.12), height=Inches(1.1),
                          font_size=Pt(60), bold=True,
                          color=self._score_color(comp_val),
                          align=PP_ALIGN.CENTER)
        self._add_textbox(slide,
                          text="分",
                          left=Inches(0.18), top=body_top + Inches(1.48),
                          width=left_w - Inches(0.12), height=Inches(0.3),
                          font_size=Pt(13),
                          color=RGBColor(0x88, 0xAA, 0xCC),
                          align=PP_ALIGN.CENTER)

        # 分隔线
        sep1 = slide.shapes.add_shape(1,
                                       Inches(0.25), body_top + Inches(1.85),
                                       left_w - Inches(0.25), Pt(1))
        sep1.fill.solid()
        sep1.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
        sep1.line.fill.background()

        # 五维分数（数字列表）
        dim_scores = [
            ("培训能力", scores.get("training")),
            ("安全意识", scores.get("safety")),
            ("工作绩效", scores.get("performance")),
            ("安全趋势", scores.get("learning")),
            ("稳  定  性", scores.get("stability")),
        ]
        y_dim = body_top + Inches(1.98)
        for label, val in dim_scores:
            val_str = f"{val:.1f}" if val is not None else "--"
            self._add_textbox(slide,
                              text=label,
                              left=Inches(0.18), top=y_dim,
                              width=Inches(1.1), height=Inches(0.38),
                              font_size=Pt(10),
                              color=RGBColor(0x88, 0xAA, 0xCC),
                              align=PP_ALIGN.LEFT)
            self._add_textbox(slide,
                              text=val_str,
                              left=Inches(1.3), top=y_dim,
                              width=Inches(0.9), height=Inches(0.38),
                              font_size=Pt(13), bold=True,
                              color=self._score_color(val),
                              align=PP_ALIGN.RIGHT)
            y_dim += Inches(0.42)

        # ═══════════════════════════════════════════════════════════════
        # 中栏：雷达图
        # ═══════════════════════════════════════════════════════════════
        mid_left  = Inches(2.65)
        mid_w     = Inches(4.8)
        radar_img = person.get("radar_image", "")
        if radar_img:
            self._insert_image(slide, radar_img,
                               left=mid_left + Inches(0.1),
                               top=body_top + Inches(0.1),
                               width=mid_w - Inches(0.2),
                               height=body_h - Inches(0.2))
        else:
            # 无雷达图时用占位
            placeholder = slide.shapes.add_shape(
                1, mid_left, body_top, mid_w, body_h
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
            placeholder.line.fill.background()
            self._add_textbox(slide, text="（雷达图加载中）",
                              left=mid_left, top=body_top + body_h / 2 - Inches(0.3),
                              width=mid_w, height=Inches(0.5),
                              font_size=Pt(12), color=COLOR_TEXT_MUTED,
                              align=PP_ALIGN.CENTER)

        # ═══════════════════════════════════════════════════════════════
        # 右栏：五维进度条 + 诊断意见
        # ═══════════════════════════════════════════════════════════════
        right_left = Inches(7.65)
        right_w    = SLIDE_WIDTH - right_left - Inches(0.2)
        right_top  = body_top + Inches(0.15)

        # 右栏白色卡片背景
        right_card = slide.shapes.add_shape(
            1, right_left - Inches(0.1), body_top,
            right_w + Inches(0.1), body_h
        )
        right_card.fill.solid()
        right_card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        right_card.line.fill.background()

        # 五维进度条
        bar_defs = [
            ("培训能力", scores.get("training"),     COLOR_PRIMARY),
            ("安全意识", scores.get("safety"),        COLOR_SUCCESS),
            ("工作绩效", scores.get("performance"),   RGBColor(0xFD, 0x7E, 0x14)),
            ("安全趋势", scores.get("learning"),      RGBColor(0x17, 0xA2, 0xB8)),
            ("稳定性",   scores.get("stability"),     COLOR_TEXT_MUTED),
        ]
        y_bar = right_top + Inches(0.05)
        bar_total_w = right_w - Inches(0.25)

        for label, val, color in bar_defs:
            val_num = val if val is not None else 0
            bar_ratio = min(max(val_num / 100.0, 0), 1.0)

            # 标签文字
            self._add_textbox(slide, text=label,
                              left=right_left, top=y_bar,
                              width=Inches(1.1), height=Inches(0.3),
                              font_size=Pt(10), color=COLOR_TEXT_MUTED,
                              align=PP_ALIGN.LEFT)

            # 分数文字
            val_str = f"{val:.0f}" if val is not None else "--"
            self._add_textbox(slide, text=val_str,
                              left=right_left + Inches(1.15), top=y_bar,
                              width=Inches(0.5), height=Inches(0.3),
                              font_size=Pt(10), bold=True,
                              color=color, align=PP_ALIGN.LEFT)

            # 进度条背景（灰色轨道）
            track_top = y_bar + Inches(0.3)
            track_h   = Pt(9)
            track = slide.shapes.add_shape(
                1, right_left, track_top, bar_total_w, track_h
            )
            track.fill.solid()
            track.fill.fore_color.rgb = RGBColor(0xE9, 0xEC, 0xEF)
            track.line.fill.background()

            # 进度条填充
            if val_num > 0:
                fill_bar = slide.shapes.add_shape(
                    1, right_left, track_top,
                    int(bar_total_w * bar_ratio), track_h
                )
                fill_bar.fill.solid()
                fill_bar.fill.fore_color.rgb = color
                fill_bar.line.fill.background()

            y_bar += Inches(0.62)

        # 分隔线
        sep_y = y_bar + Inches(0.05)
        sep2 = slide.shapes.add_shape(
            1, right_left, sep_y, right_w - Inches(0.15), Pt(1)
        )
        sep2.fill.solid()
        sep2.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        sep2.line.fill.background()

        # ── 各维度明细概览（复用画像页面概览栏数据）────────────────────
        diag_top = sep_y + Inches(0.12)
        details = person.get('details', {})
        sd = details.get('safety_details', {})
        td = details.get('training_details', {})
        pd = details.get('performance_details', {})
        ld = details.get('learning_details', {})
        std = details.get('stability_details', {})

        # 收集所有维度概览行
        overview_lines = []

        # 安全记录
        if sd:
            overview_lines.append(('⛔ 安全记录', [
                f"违规次数: {sd.get('violations', 0)}",
                f"累计扣分: {sd.get('total_deduction', 0)}",
                f"月均频次: {sd.get('avg_freq', 0)}",
                f"行为分(A): {sd.get('score_a', '--')}",
                f"严重分(B): {sd.get('score_b', '--')}",
                f"发现问题: {sd.get('as_rectifier', 0)}",
            ]))

        # 培训记录
        if td:
            overview_lines.append(('📝 培训记录', [
                f"培训总次: {td.get('total_ops', 0)}",
                f"不合格次: {td.get('fail_count', 0)}",
                f"评分: {td.get('radar_score', '--')}",
            ]))

        # 绩效
        if pd:
            overview_lines.append(('📈 工作绩效', [
                f"等级: {pd.get('display_label', '--')}",
                f"考核次数: {pd.get('count', 0)}",
            ]))

        # 安全趋势
        if ld:
            risk_map = {'SAFE': '安全', 'WATCH_LIST': '观察', 'HIGH_RISK': '高危', 'PRE_ACCIDENT': '预警'}
            risk_label = risk_map.get(ld.get('risk_level', ''), ld.get('risk_level', '--'))
            overview_lines.append(('📉 安全趋势', [
                f"风险等级: {risk_label}",
                f"本期违规: {ld.get('current_violations', 0)}",
            ]))

        # 稳定性
        if std:
            overview_lines.append(('🔄 稳定性', [
                f"{std.get('stability_label', '暂无数据')}",
            ]))

        # 渲染到 slide
        y_ov = diag_top
        for section_title, items in overview_lines:
            # 维度标题
            self._add_textbox(slide, text=section_title,
                              left=right_left, top=y_ov,
                              width=right_w, height=Inches(0.28),
                              font_size=Pt(9), bold=True,
                              color=self.theme['text_main'])
            y_ov += Inches(0.26)
            # 明细文字（一行展示，用 | 分隔）
            detail_text = '  |  '.join(items)
            self._add_textbox(slide, text=detail_text,
                              left=right_left, top=y_ov,
                              width=right_w, height=Inches(0.26),
                              font_size=Pt(8),
                              color=RGBColor(0x55, 0x55, 0x55),
                              word_wrap=True)
            y_ov += Inches(0.3)

    def _score_color(self, val) -> RGBColor:
        """根据分值返回颜色：红/橙/绿"""
        if val is None:
            return RGBColor(0xAA, 0xAA, 0xAA)
        if val < 65:
            return COLOR_DANGER
        if val < 75:
            return RGBColor(0xFD, 0x7E, 0x14)
        return COLOR_SUCCESS

    # ─────────────────────────────────────────────────────────────────────────
    # 辅助工具方法
    # ─────────────────────────────────────────────────────────────────────────

    def _add_text_only_slide(self, title: str, note: str = ""):
        """无图表时的说明页：蓝色标题栏 + 居中提示文字"""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._draw_slide_header(slide, title)
        self._add_textbox(slide,
                          text=note or f"{title}数据",
                          left=Inches(1.0),
                          top=Inches(2.5),
                          width=SLIDE_WIDTH - Inches(2.0),
                          height=Inches(2.5),
                          font_size=Pt(16),
                          color=COLOR_TEXT_MUTED,
                          align=PP_ALIGN.CENTER,
                          word_wrap=True)

    def _draw_slide_header(self, slide, title: str):
        header_height = Inches(0.8)
        bar = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, header_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme['header_bg']
        bar.line.fill.background()

        self._add_textbox(slide, text=title,
                          left=Inches(0.3), top=Inches(0.1),
                          width=SLIDE_WIDTH - Inches(0.6), height=header_height - Inches(0.1),
                          font_size=Pt(20), bold=True,
                          color=self.theme['header_text'])

    def _add_textbox(self, slide, text: str, left, top, width, height,
                     font_size=Pt(14), bold=False,
                     color=None, align=PP_ALIGN.LEFT, word_wrap=False):
        """添加文本框"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return txBox

    def _insert_image(self, slide, img_b64: str, left, top, width, height):
        """将 base64 图片插入幻灯片，保持原始比例居中放置"""
        if not img_b64:
            return
        # 去掉 data:image/xxx;base64, 前缀
        if "," in img_b64:
            img_b64 = img_b64.split(",", 1)[1]
        try:
            img_bytes = base64.b64decode(img_b64)
            img_stream = io.BytesIO(img_bytes)

            # 读取原始图片尺寸，按比例缩放到容器内居中
            from PIL import Image as PILImage
            pil_img = PILImage.open(io.BytesIO(img_bytes))
            img_w, img_h = pil_img.size
            pil_img.close()

            # 将容器尺寸转为数值（EMU）
            container_w = int(width)
            container_h = int(height)

            # 计算等比缩放
            scale_w = container_w / img_w
            scale_h = container_h / img_h
            scale = min(scale_w, scale_h)  # fit 模式

            actual_w = int(img_w * scale)
            actual_h = int(img_h * scale)

            # 居中偏移
            offset_left = int(left) + (container_w - actual_w) // 2
            offset_top  = int(top)  + (container_h - actual_h) // 2

            img_stream.seek(0)
            slide.shapes.add_picture(img_stream, offset_left, offset_top, actual_w, actual_h)
        except Exception as e:
            # 图片解码失败：插入占位文字，不中断整体PPT生成
            self._add_textbox(slide,
                              text=f"[图表加载失败: {e}]",
                              left=left, top=top, width=width, height=height,
                              font_size=Pt(10), color=COLOR_TEXT_MUTED,
                              align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# AI 文字服务（含缓存 + 规则降级）
# ─────────────────────────────────────────────────────────────────────────────

class AITextService:
    """
    为关键人员生成诊断文字
    优先从 ppt_export_cache 取缓存，无则调AI，AI失败则降级规则拼接
    """

    CACHE_TTL_DAYS = 30

    def get_summary(self, emp_no: str, start_date: str, end_date: str,
                    profile_data: dict) -> str:
        """
        获取诊断文字（带缓存）
        返回：str（永不抛出异常）
        """
        cache_key = self._make_cache_key(emp_no, start_date, end_date, profile_data)

        # 1. 查缓存
        cached = self._get_cache(cache_key)
        if cached:
            return cached

        # 2. 尝试AI
        try:
            summary = self._call_ai(emp_no, profile_data)
            self._set_cache(cache_key, emp_no, start_date, end_date, summary,
                            is_ai=True, tokens=self._last_tokens)
            return summary
        except Exception as e:
            print(f"[PPT AI] AI调用失败，降级规则文字: {e}")

        # 3. 降级规则拼接
        summary = self._rule_summary(profile_data)
        self._set_cache(cache_key, emp_no, start_date, end_date, summary,
                        is_ai=False, tokens=0)
        return summary

    # ── 缓存操作 ──────────────────────────────────────────────────────────────

    def _make_cache_key(self, emp_no: str, start_date: str, end_date: str,
                        profile_data: dict) -> str:
        scores = profile_data.get("scores", {})
        safety = profile_data.get("safety_details", {})
        fingerprint = hashlib.md5(
            json.dumps({
                "scores": scores,
                "violations": safety.get("violations"),
            }, sort_keys=True).encode()
        ).hexdigest()[:12]
        raw = f"{emp_no}|{start_date}|{end_date}|{fingerprint}"
        return hashlib.md5(raw.encode()).hexdigest()

    def _get_cache(self, cache_key: str):
        try:
            from models.database import get_db
            conn = get_db()
            cur = conn.cursor()
            cur.execute(
                "SELECT ai_summary FROM ppt_export_cache "
                "WHERE cache_key = %s AND expires_at > NOW()",
                (cache_key,)
            )
            row = cur.fetchone()
            return row["ai_summary"] if row else None
        except Exception:
            return None

    def _set_cache(self, cache_key: str, emp_no: str, start_date: str,
                   end_date: str, summary: str,
                   is_ai: bool = False, tokens: int = 0):
        try:
            from models.database import get_db
            conn = get_db()
            cur = conn.cursor()
            expires_at = datetime.now() + timedelta(days=self.CACHE_TTL_DAYS)
            cur.execute(
                """
                INSERT INTO ppt_export_cache
                    (cache_key, emp_no, start_date, end_date, ai_summary,
                     is_ai_generated, tokens_used, expires_at)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
                ON DUPLICATE KEY UPDATE
                    ai_summary = VALUES(ai_summary),
                    is_ai_generated = VALUES(is_ai_generated),
                    tokens_used = VALUES(tokens_used),
                    expires_at = VALUES(expires_at)
                """,
                (cache_key, emp_no, start_date, end_date, summary,
                 1 if is_ai else 0, tokens, expires_at)
            )
            conn.commit()
        except Exception as e:
            print(f"[PPT AI] 写入缓存失败（非致命）: {e}")

    # ── AI 调用 ───────────────────────────────────────────────────────────────

    _last_tokens: int = 0

    def _call_ai(self, emp_no: str, profile_data: dict) -> str:
        """调用已有AI基础设施生成2-3句诊断文字，复用AIDiagnosisService的配置"""
        import httpx
        from services.ai_diagnosis_service import AIDiagnosisService

        config = AIDiagnosisService._get_ai_config()
        if not config or not config.get("api_key"):
            raise RuntimeError("AI未配置")

        scores   = profile_data.get("scores", {})
        safety   = profile_data.get("safety_details", {})
        training = profile_data.get("training_details", {})
        stability= profile_data.get("stability_details", {})
        learning = profile_data.get("learning_details", {})

        prompt = (
            f"请用2-3句话为以下员工给出简洁的管理建议（不超过120字）：\n"
            f"综合分：{scores.get('comprehensive', '--')}分，"
            f"培训能力：{scores.get('training', '--')}分，"
            f"安全意识：{scores.get('safety', '--')}分，"
            f"绩效：{scores.get('performance', '--')}分，"
            f"安全趋势：{scores.get('learning', '--')}分，"
            f"稳定性：{scores.get('stability', '--')}分。\n"
            f"安全告警：{safety.get('alert_tag', '无')}，"
            f"违规次数：{safety.get('violations', 0)}次，"
            f"培训告警：{training.get('alert_tag', '无')}，"
            f"稳定度告警：{stability.get('alert_tag', '无')}，"
            f"安全趋势提示：{learning.get('alert_tag', '无')}。\n"
            "请给出针对性管理建议，不要空话套话。"
        )

        provider_type = config.get("provider_type", "openai")
        api_key   = config["api_key"]
        model     = config.get("model", "gpt-3.5-turbo")
        base_url  = config.get("base_url", "https://api.openai.com/v1")
        timeout   = config.get("timeout", 30)
        max_tokens= min(config.get("max_tokens", 200), 300)
        temperature = config.get("temperature", 0.7)
        extra_headers = config.get("extra_headers") or {}

        headers = {"Content-Type": "application/json",
                   "Authorization": f"Bearer {api_key}"}
        if provider_type == "anthropic":
            headers["x-api-key"] = api_key
            del headers["Authorization"]
        elif provider_type == "gemini":
            headers.pop("Authorization", None)
        headers.update(extra_headers)

        if provider_type == "gemini":
            payload = {"contents": [{"parts": [{"text": prompt}]}],
                       "generationConfig": {"temperature": temperature,
                                            "maxOutputTokens": max_tokens}}
            endpoint = f"{base_url}/models/{model}:generateContent?key={api_key}"
        elif provider_type == "anthropic":
            payload = {"model": model,
                       "messages": [{"role": "user", "content": prompt}],
                       "max_tokens": max_tokens}
            endpoint = f"{base_url}/messages"
        else:
            payload = {"model": model,
                       "messages": [{"role": "user", "content": prompt}],
                       "max_tokens": max_tokens,
                       "temperature": temperature}
            endpoint = f"{base_url}/chat/completions"

        with httpx.Client(timeout=timeout) as client:
            resp = client.post(endpoint, headers=headers, json=payload)
            resp.raise_for_status()
            data = resp.json()

        if provider_type == "gemini":
            candidates = data.get("candidates", [{}])
            parts = candidates[0].get("content", {}).get("parts", [{}]) if candidates else [{}]
            content = parts[0].get("text", "") if parts else ""
            tokens = data.get("usageMetadata", {}).get("totalTokenCount", 0)
        elif provider_type == "anthropic":
            content = data.get("content", [{}])[0].get("text", "")
            usage = data.get("usage", {})
            tokens = usage.get("input_tokens", 0) + usage.get("output_tokens", 0)
        else:
            content = data["choices"][0]["message"]["content"].strip()
            tokens = data.get("usage", {}).get("total_tokens", 0)

        self._last_tokens = tokens

        # 记录使用量（复用现有日志）
        try:
            AIDiagnosisService._log_usage(
                config.get("id"), config.get("name", "Unknown"),
                model, tokens, True, None
            )
        except Exception:
            pass

        return content.strip()


    # ── 规则降级 ──────────────────────────────────────────────────────────────

    def _rule_summary(self, profile_data: dict) -> str:
        """不依赖 AI 的规则拼接，永不失败"""
        scores   = profile_data.get("scores", {})
        safety   = profile_data.get("safety_details", {})
        training = profile_data.get("training_details", {})
        stability= profile_data.get("stability_details", {})
        learning = profile_data.get("learning_details", {})

        issues = []
        suggestions = []

        comp = scores.get("comprehensive")
        if comp is not None and comp < 75:
            issues.append(f"综合评分偏低（{comp:.1f}分）")

        safe_score = scores.get("safety")
        violations = safety.get("violations", 0)
        if violations and violations > 0:
            issues.append(f"安全违规{violations}次")
            suggestions.append("加强安全日常督导")
        if safe_score is not None and safe_score < 70:
            issues.append(f"安全意识评分{safe_score:.1f}分")

        train_score = scores.get("training")
        fail_count  = training.get("fail_count", 0)
        if fail_count and fail_count > 0:
            issues.append(f"培训失格{fail_count}次")
            suggestions.append("安排重点培训辅导")
        if train_score is not None and train_score < 70:
            issues.append(f"培训能力评分{train_score:.1f}分")

        stab_alert = stability.get("alert_tag", "")
        if stab_alert and "⚠" in stab_alert:
            issues.append(stab_alert.replace("⚠️", "").replace("⚠", "").strip())

        learn_alert = learning.get("alert_tag", "")
        if learn_alert and ("⚠" in learn_alert or "危" in learn_alert):
            issues.append("安全趋势恶化")
            suggestions.append("关注近期行为变化")

        if not issues:
            return "综合评分低于关键人员阈值，建议持续关注各维度表现。"

        summary = "主要问题：" + "；".join(issues) + "。"
        if suggestions:
            summary += "建议：" + "；".join(suggestions) + "。"
        return summary
